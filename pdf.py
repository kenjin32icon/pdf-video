import tkinter as tk
from tkinter import filedialog, ttk
import PyPDF2
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract
import os
from gtts import gTTS
import threading
from pathlib import Path
import cv2
import numpy as np

class DocumentToAudio:
    def __init__(self):
        self.window = tk.Tk()
        self.window.title("Document to Audio Converter")
        self.window.geometry("800x600")
        
        # Create main frame
        self.main_frame = ttk.Frame(self.window, padding="10")
        self.main_frame.grid(row=0, column=0, sticky=(tk.W, tk.E, tk.N, tk.S))
        
        # Create UI elements
        self.create_widgets()
        
        # Initialize variables
        self.file_path = None
        self.extracted_text = ""
        
    def create_widgets(self):
        # File selection
        ttk.Button(self.main_frame, text="Select File", command=self.select_file).grid(row=0, column=0, pady=5)
        self.file_label = ttk.Label(self.main_frame, text="No file selected")
        self.file_label.grid(row=0, column=1, pady=5)
        
        # Text display
        self.text_display = tk.Text(self.main_frame, height=15, width=70)
        self.text_display.grid(row=1, column=0, columnspan=2, pady=10)
        
        # Progress bar
        self.progress = ttk.Progressbar(self.main_frame, length=300, mode='determinate')
        self.progress.grid(row=2, column=0, columnspan=2, pady=5)
        
        # Convert button
        ttk.Button(self.main_frame, text="Convert to Audio", command=self.convert_to_audio).grid(row=3, column=0, columnspan=2, pady=5)
        
        # Status label
        self.status_label = ttk.Label(self.main_frame, text="")
        self.status_label.grid(row=4, column=0, columnspan=2, pady=5)
    
    def select_file(self):
        filetypes = (
            ('PDF files', '*.pdf'),
            ('Word files', '*.docx'),
            ('PowerPoint files', '*.pptx'),
            ('Text files', '*.txt'),
            ('Image files', '*.png;*.jpg;*.jpeg')
        )
        
        self.file_path = filedialog.askopenfilename(filetypes=filetypes)
        if self.file_path:
            self.file_label.config(text=os.path.basename(self.file_path))
            self.extract_text()
    
    def extract_text(self):
        self.status_label.config(text="Extracting text...")
        self.progress['value'] = 0
        
        file_extension = os.path.splitext(self.file_path)[1].lower()
        
        try:
            if file_extension == '.pdf':
                self.extract_from_pdf()
            elif file_extension == '.docx':
                self.extract_from_docx()
            elif file_extension == '.pptx':
                self.extract_from_pptx()
            elif file_extension == '.txt':
                self.extract_from_txt()
            elif file_extension in ['.png', '.jpg', '.jpeg']:
                self.extract_from_image()
            
            self.text_display.delete(1.0, tk.END)
            self.text_display.insert(tk.END, self.extracted_text)
            self.status_label.config(text="Text extracted successfully!")
            self.progress['value'] = 100
            
        except Exception as e:
            self.status_label.config(text=f"Error: {str(e)}")
    
    def extract_from_pdf(self):
        with open(self.file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            text = ""
            total_pages = len(pdf_reader.pages)
            
            for i, page in enumerate(pdf_reader.pages):
                text += page.extract_text() + "\n"
                self.progress['value'] = ((i + 1) / total_pages) * 100
                self.window.update_idletasks()
            
            self.extracted_text = text
    
    def extract_from_docx(self):
        doc = Document(self.file_path)
        self.extracted_text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
    
    def extract_from_pptx(self):
        prs = Presentation(self.file_path)
        text = ""
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        
        self.extracted_text = text
    
    def extract_from_txt(self):
        with open(self.file_path, 'r', encoding='utf-8') as file:
            self.extracted_text = file.read()
    
    def extract_from_image(self):
        image = cv2.imread(self.file_path)
        gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
        self.extracted_text = pytesseract.image_to_string(gray)
    
    def convert_to_audio(self):
        if not self.extracted_text:
            self.status_label.config(text="Please extract text first!")
            return
        
        self.status_label.config(text="Converting to audio...")
        threading.Thread(target=self._convert_to_audio_thread).start()
    
    def _convert_to_audio_thread(self):
        try:
            output_path = os.path.splitext(self.file_path)[0] + '.mp3'
            tts = gTTS(text=self.extracted_text, lang='en')
            tts.save(output_path)
            self.status_label.config(text=f"Audio saved as: {os.path.basename(output_path)}")
        except Exception as e:
            self.status_label.config(text=f"Error converting to audio: {str(e)}")
    
    def run(self):
        self.window.mainloop()

if __name__ == "__main__":
    app = DocumentToAudio()
    app.run()